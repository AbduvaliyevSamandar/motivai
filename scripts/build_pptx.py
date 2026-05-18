# -*- coding: utf-8 -*-
"""MotivAI thesis defence deck — text-first version.

Each slide carries a short headline, 3–5 concise bullets, and a small
native-drawn diagram (shapes + arrows). No large PNG imports — speech-
friendly, easy to scan at a glance, suitable for live presentation.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "MotivAI — Diplom Taqdimoti.pptx"

# Brand palette
PRIMARY   = RGBColor(0x4F, 0x46, 0xE5)
SECONDARY = RGBColor(0x7C, 0x3A, 0xED)
ACCENT    = RGBColor(0xF5, 0x9E, 0x0B)
SUCCESS   = RGBColor(0x10, 0xB9, 0x81)
DANGER    = RGBColor(0xEF, 0x44, 0x44)
SKY       = RGBColor(0x0E, 0xA5, 0xE9)
INK       = RGBColor(0x0F, 0x17, 0x2A)
TXT       = RGBColor(0x1F, 0x29, 0x37)
SUB       = RGBColor(0x6B, 0x72, 0x80)
BG        = RGBColor(0xFA, 0xFA, 0xFB)
TINT      = RGBColor(0xEE, 0xF2, 0xFF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)


# ── helpers ────────────────────────────────────────────────────────────
def setup(prs):
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def text(slide, x, y, w, h, body, *, size=18, bold=False, color=TXT,
         align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = body
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def bullets(slide, x, y, w, h, items, *, size=14, color=TXT,
            bullet_color=PRIMARY, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        p.line_spacing = line_spacing
        b = p.add_run()
        b.text = "•  "
        b.font.color.rgb = bullet_color
        b.font.size = Pt(size + 2)
        b.font.bold = True
        b.font.name = "Calibri"
        body = p.add_run()
        body.text = line
        body.font.size = Pt(size)
        body.font.color.rgb = color
        body.font.name = "Calibri"
    return tb


def rect(slide, x, y, w, h, *, fill=PRIMARY, line=None, line_w=1.2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line; shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def rounded(slide, x, y, w, h, *, fill=TINT, line=None, corner=0.10, line_w=1.2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = corner
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line; shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def circle(slide, x, y, d, *, fill=PRIMARY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def arrow(slide, x1, y1, x2, y2, *, color=SUB, weight=1.4):
    line = slide.shapes.add_connector(2, x1, y1, x2, y2)  # straight
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def header(slide, title, subtitle=None):
    # left accent strip
    rect(slide, Inches(0), Inches(0), Inches(0.10), Inches(7.5), fill=PRIMARY)
    text(slide, Inches(0.55), Inches(0.30), Inches(12.2), Inches(0.65),
         title, size=28, bold=True, color=INK)
    if subtitle:
        text(slide, Inches(0.55), Inches(1.00), Inches(12.2), Inches(0.40),
             subtitle, size=14, color=SUB, italic=True)


def footer(slide, page, total=20):
    text(slide, Inches(0.55), Inches(7.05), Inches(8.0), Inches(0.30),
         "MotivAI · Abduvaliyev Samandar · TATU 2026",
         size=9, color=SUB)
    text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.30),
         f"{page} / {total}", size=9, color=SUB, align=PP_ALIGN.RIGHT)


def chip(slide, x, y, w, h, label, *, fill=PRIMARY, txt_color=WHITE, size=11):
    rounded(slide, x, y, w, h, fill=fill, corner=0.30)
    text(slide, x, y + Inches(0.05), w, h - Inches(0.05),
         label, size=size, bold=True, color=txt_color, align=PP_ALIGN.CENTER)


# ── slides ─────────────────────────────────────────────────────────────
def s01_title(prs):
    s = blank(prs)
    rect(s, 0, 0, prs.slide_width, Inches(4.2), fill=PRIMARY)
    rect(s, 0, Inches(4.2), prs.slide_width, Inches(3.3), fill=BG)

    # Logo emblem
    rounded(s, Inches(5.85), Inches(0.8), Inches(1.65), Inches(1.65),
            fill=WHITE, corner=0.20)
    text(s, Inches(5.85), Inches(0.95), Inches(1.65), Inches(1.4),
         "M↑", size=64, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    text(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.7),
         "MotivAI", size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, Inches(0.5), Inches(3.45), Inches(12.3), Inches(0.5),
         "Sun'iy intellekt asosida talabalar motivatsiyasini boshqaruvchi mobil platforma",
         size=16, color=RGBColor(0xE0, 0xE7, 0xFF), italic=True, align=PP_ALIGN.CENTER)

    text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.4),
         "DIPLOM LOYIHASI HIMOYASI",
         size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    # contact card
    rounded(s, Inches(3.5), Inches(5.25), Inches(6.3), Inches(1.5),
            fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB), corner=0.10)
    text(s, Inches(3.5), Inches(5.35), Inches(6.3), Inches(0.35),
         "Bitiruvchi", size=10, color=SUB, align=PP_ALIGN.CENTER)
    text(s, Inches(3.5), Inches(5.65), Inches(6.3), Inches(0.4),
         "Abduvaliyev Samandar Qobil o'g'li",
         size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
    text(s, Inches(3.5), Inches(6.1), Inches(6.3), Inches(0.35),
         "Rahbar: Sodiqov V.S.   ·   Yo'nalish: 60610500", size=11,
         color=SUB, align=PP_ALIGN.CENTER)
    text(s, Inches(3.5), Inches(6.4), Inches(6.3), Inches(0.35),
         "TATU · Toshkent · 2026", size=10,
         color=SUB, italic=True, align=PP_ALIGN.CENTER)


def s02_problem(prs):
    s = blank(prs)
    header(s, "Muammoning dolzarbligi",
           "Talabalarda motivatsiya pasayishi — global va mahalliy kontekst")

    # 3 stat circles on top
    stats = [
        ("53%", "Motivatsion qiyinchiliklar (UNESCO 2023)",     DANGER),
        ("48%", "Ikkinchi semestrgacha pasayish (TATU 2022)",   ACCENT),
        ("38%", "Ta'lim maqsadlaridan uzilish xavfi",            SECONDARY),
    ]
    for i, (n, lbl, col) in enumerate(stats):
        x = Inches(0.8 + i * 4.15)
        circle(s, x, Inches(1.8), Inches(1.7), fill=col)
        text(s, x, Inches(2.0), Inches(1.7), Inches(1.0),
             n, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, x, Inches(3.65), Inches(1.7), Inches(0.9),
             lbl, size=11, color=TXT, align=PP_ALIGN.CENTER)

    # Bottom statement
    rounded(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.7),
            fill=TINT, corner=0.06)
    text(s, Inches(0.95), Inches(5.15), Inches(11.5), Inches(0.4),
         "Yechim: shaxsiylashtirilgan motivatsion platforma",
         size=16, bold=True, color=PRIMARY)
    bullets(s, Inches(0.95), Inches(5.55), Inches(11.5), Inches(1.1), [
        "AI orqali har kuni shaxsiy vazifa rejasi tuziladi",
        "Gamifikatsiya (streak, XP, daraja) — barqaror ishtirok",
        "O'zbek tilida, mahalliy kontekstga moslashgan",
    ], size=13)
    footer(s, 2)


def s03_goal_tasks(prs):
    s = blank(prs)
    header(s, "Maqsad va vazifalar")

    # Goal card
    rounded(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(1.4),
            fill=PRIMARY, corner=0.06)
    text(s, Inches(0.85), Inches(1.65), Inches(11.6), Inches(0.35),
         "MAQSAD", size=10, bold=True, color=RGBColor(0xE0, 0xE7, 0xFF))
    text(s, Inches(0.85), Inches(2.0), Inches(11.6), Inches(0.85),
         "Sun'iy intellekt orqali talabaning shaxsiy motivatsiya rejasini har kuni "
         "avtomatik tuzib beruvchi va gamifikatsiya orqali barqaror ishtirokni "
         "ta'minlovchi cross-platform mobil ilovani ishlab chiqish.",
         size=13, color=WHITE)

    # 6 tasks as numbered chips
    text(s, Inches(0.6), Inches(3.2), Inches(12.1), Inches(0.4),
         "ASOSIY VAZIFALAR", size=11, bold=True, color=PRIMARY)
    tasks = [
        "Tavsiya tizimlari va ta'limda AI yo'nalishlarining nazariy tahlili",
        "MVF (Motivatsion Qiymat Funksiyasi) matematik modelini shakllantirish",
        "Flutter (Dart) bilan iOS+Android uchun mobil ilovani ishlab chiqish",
        "FastAPI orqali modulli RESTful API arxitekturasini yaratish",
        "MongoDB Atlas ma'lumotlar bazasini loyihalash va indekslash",
        "Ko'p providerli AI fallback chain (OpenAI/Gemini/Groq) integratsiyasi",
    ]
    for i, t in enumerate(tasks):
        col, row = i % 2, i // 2
        x = Inches(0.6 + col * 6.15)
        y = Inches(3.7 + row * 1.0)
        circle(s, x, y, Inches(0.55), fill=PRIMARY)
        text(s, x, y + Inches(0.08), Inches(0.55), Inches(0.4),
             f"{i+1}", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.75), y + Inches(0.05), Inches(5.4), Inches(0.85),
             t, size=12, color=TXT)
    footer(s, 3)


def s04_architecture_mini(prs):
    s = blank(prs)
    header(s, "Tizim arxitekturasi", "Uch qatlamli klassik dizayn")

    # Mini stacked architecture
    layers = [
        ("MIJOZ QATLAMI · Flutter (Dart)",
         "Dashboard · AI Chat · Leaderboard · Progress · Profile",
         PRIMARY),
        ("BIZNES MANTIQ · FastAPI (Python)",
         "Auth · Tasks · Chat · AI · MVF · Gamification",
         SECONDARY),
        ("MA'LUMOTLAR · MongoDB Atlas",
         "users · tasks · progress · chat_sessions · motivation_plans",
         SUCCESS),
    ]
    for i, (title, sub, col) in enumerate(layers):
        y = Inches(1.85 + i * 1.55)
        rounded(s, Inches(1.5), y, Inches(10.3), Inches(1.3),
                fill=WHITE, line=col, corner=0.05, line_w=2.0)
        rect(s, Inches(1.5), y, Inches(0.15), Inches(1.3), fill=col)
        text(s, Inches(1.85), y + Inches(0.18), Inches(9.5), Inches(0.45),
             title, size=15, bold=True, color=col)
        text(s, Inches(1.85), y + Inches(0.75), Inches(9.5), Inches(0.5),
             sub, size=12, color=TXT)
        if i < 2:
            arrow(s, Inches(6.6), y + Inches(1.3), Inches(6.6), y + Inches(1.55),
                  color=SUB, weight=2)

    # Labels on arrows
    text(s, Inches(7.0), Inches(3.15), Inches(2.0), Inches(0.3),
         "HTTPS / JWT", size=10, color=SUB, italic=True)
    text(s, Inches(7.0), Inches(4.7), Inches(2.0), Inches(0.3),
         "Motor (async)", size=10, color=SUB, italic=True)

    # Side note
    text(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
         "Har qatlam mustaqil — alohida testlanadi, monitoring qilinadi, miqyoslanadi",
         size=12, color=SUB, italic=True, align=PP_ALIGN.CENTER)
    footer(s, 4)


def s05_techstack(prs):
    s = blank(prs)
    header(s, "Texnologik stack", "Zamonaviy, asoslangan tanlovlar")

    stack = [
        ("Mobile UI",     "Flutter 3.16",  "Dart · hot reload · 1 codebase iOS+Android",  PRIMARY),
        ("Backend",       "FastAPI 0.111", "ASGI · Pydantic v2 · auto-docs (/docs)",      SECONDARY),
        ("Database",      "MongoDB Atlas", "M0 Singapore · Motor async · 5 collections",  SUCCESS),
        ("AI (multi-LLM)","3 providers",   "OpenAI gpt-4o-mini → Gemini → Groq Llama",    ACCENT),
        ("Hosting",       "Render.com",    "Auto HTTPS · CI/CD on git push",              SKY),
        ("Auth",          "JWT + bcrypt",  "HMAC-SHA256 · work factor 12",                DANGER),
    ]
    for i, (cat, name, desc, col) in enumerate(stack):
        c, r = i % 3, i // 3
        x = Inches(0.6 + c * 4.15)
        y = Inches(1.85 + r * 2.4)
        rounded(s, x, y, Inches(3.95), Inches(2.1),
                fill=WHITE, line=col, corner=0.06, line_w=2.0)
        text(s, x + Inches(0.25), y + Inches(0.20), Inches(3.6), Inches(0.4),
             cat, size=11, color=SUB)
        text(s, x + Inches(0.25), y + Inches(0.60), Inches(3.6), Inches(0.6),
             name, size=20, bold=True, color=col)
        text(s, x + Inches(0.25), y + Inches(1.25), Inches(3.6), Inches(0.75),
             desc, size=11, color=TXT)
    footer(s, 5)


def s06_archetypes(prs):
    s = blank(prs)
    header(s, "Foydalanuvchi arxetiplari",
           "K-means klasterlash (K=5, silhouette = 0.62)")

    archetypes = [
        ("1", "Boshlang'ich", "Beginner",     "Vazifa = 0",          SUCCESS),
        ("2", "Tadqiqotchi",   "Explorer",     "Haftalik < 5",        SKY),
        ("3", "Izchil",        "Consistent",   "Streak ≥ 3",          SECONDARY),
        ("4", "Muvaffaqiyatli","Achiever",     "Haftalik > 5",        ACCENT),
        ("5", "Chempion",      "Champion",     "Streak ≥ 14",         DANGER),
    ]
    w = Inches(2.40)
    pad = Inches(0.12)
    x0 = Inches(0.6)
    for i, (n, name, en, crit, col) in enumerate(archetypes):
        x = x0 + (w + pad) * i
        # number circle
        circle(s, x + Inches(0.85), Inches(2.0), Inches(0.7), fill=col)
        text(s, x + Inches(0.85), Inches(2.05), Inches(0.7), Inches(0.6),
             n, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # name
        text(s, x, Inches(2.9), w, Inches(0.4),
             name, size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
        text(s, x, Inches(3.3), w, Inches(0.3),
             en, size=10, color=SUB, italic=True, align=PP_ALIGN.CENTER)
        # criterion chip
        chip(s, x + Inches(0.3), Inches(3.85), w - Inches(0.6), Inches(0.45),
             crit, fill=col, size=11)
        if i < len(archetypes) - 1:
            arrow(s, x + w, Inches(2.35),
                  x + w + pad, Inches(2.35), color=col, weight=2.0)

    # Bottom: transition note
    rounded(s, Inches(0.6), Inches(4.95), Inches(12.1), Inches(1.7),
            fill=TINT, corner=0.06)
    text(s, Inches(0.95), Inches(5.10), Inches(11.5), Inches(0.4),
         "Har arxetip uchun maxsus motivatsion strategiya",
         size=14, bold=True, color=PRIMARY)
    bullets(s, Inches(0.95), Inches(5.55), Inches(11.5), Inches(1.0), [
        "Boshlang'ich → oson vazifalar, katta XP mukofotlar, onboarding",
        "Chempion → ekspert qiyinligi, global reyting, 4× XP multiplikator",
        "Arxetip har 6 soatda qayta hisoblab boriladi — dinamik moslashish",
    ], size=12)
    footer(s, 6)


def s07_mvf(prs):
    s = blank(prs)
    header(s, "MVF — Motivatsion Qiymat Funksiyasi",
           "To'rt komponentli gibrid model · NDCG@5 = 0.78")

    # Formula at top
    rounded(s, Inches(1.2), Inches(1.75), Inches(10.9), Inches(0.85),
            fill=PRIMARY, corner=0.10)
    text(s, Inches(1.2), Inches(1.85), Inches(10.9), Inches(0.65),
         "MVF(u, t, C)  =  0.25·CS  +  0.25·CF  +  0.35·DM  +  0.15·TS",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 4 components
    comps = [
        ("CS",  "Kontent o'xshashlik",  "0.25", "Kosinusli o'xshashlik · 9D vektor", PRIMARY),
        ("CF",  "Kollaborativ filtr",   "0.25", "K-NN (K=20) · Pearson korrelyatsiya", SECONDARY),
        ("DM",  "Qiyinlilik mosligi",   "0.35", "Gauss · δ=2 daraja · Flow nazariyasi", ACCENT),
        ("TS",  "Vaqtinchalik kontekst","0.15", "Logistik regressiya · CARS paradigmasi", SUCCESS),
    ]
    for i, (sym, name, w, sub, col) in enumerate(comps):
        x = Inches(0.6 + i * 3.1)
        y = Inches(3.0)
        rounded(s, x, y, Inches(2.9), Inches(2.7),
                fill=WHITE, line=col, corner=0.06, line_w=2.0)
        circle(s, x + Inches(1.05), y + Inches(0.20), Inches(0.8), fill=col)
        text(s, x + Inches(1.05), y + Inches(0.30), Inches(0.8), Inches(0.6),
             sym, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, x, y + Inches(1.15), Inches(2.9), Inches(0.45),
             name, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
        text(s, x, y + Inches(1.55), Inches(2.9), Inches(0.35),
             f"og'irlik {w}", size=10, color=SUB, italic=True, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.2), y + Inches(2.0), Inches(2.5), Inches(0.65),
             sub, size=10, color=TXT, align=PP_ALIGN.CENTER)

    # Bottom theory + result
    text(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.4),
         "Nazariy asos: Self-Determination Theory (SDT) + Flow Theory (Csikszentmihalyi)",
         size=11, color=SUB, italic=True, align=PP_ALIGN.CENTER)
    text(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.4),
         "Natija: kunlik K = 5 ta eng mos vazifa",
         size=13, bold=True, color=SUCCESS, align=PP_ALIGN.CENTER)
    footer(s, 7)


def s08_user_flow(prs):
    s = blank(prs)
    header(s, "Foydalanuvchi tajriba tsikli", "6 bosqichli motivatsion zanjir")

    steps = [
        ("1", "Vazifa\ntanlash",     "MVF kunlik 5 ta",           PRIMARY),
        ("2", "Bajarish",            "Timer + UI",                SECONDARY),
        ("3", "XP mukofot",          "Qiyinlik × Streak",         ACCENT),
        ("4", "Daraja\no'sishi",     "XP → keyingi level",        SUCCESS),
        ("5", "Yutuq\nnishonlari",   "8 kategoriya",              DANGER),
        ("6", "Reyting va\nraqobat", "Global + haftalik",         SKY),
    ]
    n = len(steps)
    w = Inches(1.85)
    pad = Inches(0.10)
    x0 = Inches(0.6)
    y = Inches(2.4)
    for i, (num, title, desc, col) in enumerate(steps):
        x = x0 + (w + pad) * i
        # number badge
        circle(s, x + Inches(0.55), y, Inches(0.75), fill=col)
        text(s, x + Inches(0.55), y + Inches(0.05), Inches(0.75), Inches(0.65),
             num, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # title card
        rounded(s, x, y + Inches(1.0), w, Inches(2.1),
                fill=WHITE, line=col, corner=0.06, line_w=2.0)
        text(s, x, y + Inches(1.15), w, Inches(0.8),
             title, size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        text(s, x, y + Inches(2.15), w, Inches(0.7),
             desc, size=10, color=TXT, align=PP_ALIGN.CENTER, italic=True)
        if i < n - 1:
            arrow(s, x + w, y + Inches(0.40), x + w + pad, y + Inches(0.40),
                  color=col, weight=2.0)

    # Bottom note
    text(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
         "Tsikl yopiq — 6-bosqich avtomatik 1-bosqichga qaytaradi (har kun)",
         size=12, color=SUB, italic=True, align=PP_ALIGN.CENTER)
    footer(s, 8)


def s09_ai_chain(prs):
    s = blank(prs)
    header(s, "AI Chat — ko'p providerli fallback chain",
           "Eng katta innovatsion yechim")

    # 4 providers horizontal
    providers = [
        ("OpenAI",   "gpt-4o-mini",        "Birinchi navbat\n1,8 sek o'rtacha",      DANGER),
        ("Google",   "Gemini 2.0 Flash",   "Quota tugasa\n1500/kun bepul",            ACCENT),
        ("Groq",     "Llama 3.3 70B",      "Ikkinchi quota\n30/min · juda tez",       SKY),
        ("Fallback", "Rule-based",         "Hammasi tugasa\noffline ishlaydi",        SUCCESS),
    ]
    w = Inches(2.7)
    pad = Inches(0.40)
    x0 = Inches(0.7)
    y = Inches(2.05)
    for i, (vendor, model, role, col) in enumerate(providers):
        x = x0 + (w + pad) * i
        rounded(s, x, y, w, Inches(3.4), fill=WHITE, line=col, corner=0.08, line_w=2.4)
        rect(s, x, y, w, Inches(0.50), fill=col)
        text(s, x, y + Inches(0.1), w, Inches(0.4),
             vendor, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, x, y + Inches(0.7), w, Inches(0.45),
             model, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.2), y + Inches(1.5), w - Inches(0.4), Inches(1.0),
             role, size=11, color=TXT, align=PP_ALIGN.CENTER)
        circle(s, x + w/2 - Inches(0.3), y + Inches(2.7), Inches(0.6),
               fill=col)
        text(s, x + w/2 - Inches(0.3), y + Inches(2.78), Inches(0.6), Inches(0.5),
             str(i+1), size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(providers) - 1:
            arrow(s, x + w, y + Inches(1.7), x + w + pad, y + Inches(1.7),
                  color=SUB, weight=2.5)

    # Bottom result
    rounded(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.95),
            fill=TINT, corner=0.08)
    text(s, Inches(0.95), Inches(5.97), Inches(11.5), Inches(0.45),
         "99.6% uptime kafolat",
         size=15, bold=True, color=SUCCESS)
    text(s, Inches(0.95), Inches(6.35), Inches(11.5), Inches(0.4),
         "OpenAI kvota cheklovlariga bog'liqlik 90%+ qisqartirildi",
         size=11, color=TXT, italic=True)
    footer(s, 9)


def s10_database(prs):
    s = blank(prs)
    header(s, "MongoDB Atlas — 5 kolleksiya",
           "Embedded + referenced gibrid yondashuv")

    cols = [
        ("users",            "11 field",  "email idx · embedded badges/preferences",    PRIMARY),
        ("tasks",            "9 field",   "category · difficulty · xp_reward",          SECONDARY),
        ("progress",         "7 field",   "{user_id, completed_at} compound idx",       SUCCESS),
        ("chat_sessions",    "5 field",   "embedded messages (≤100 / hujjat)",          ACCENT),
        ("motivation_plans", "7 field",   "AI generatsiya natijalari · TTL",            DANGER),
    ]
    for i, (name, fields, desc, col) in enumerate(cols):
        y = Inches(1.85 + i * 1.0)
        rounded(s, Inches(0.6), y, Inches(12.1), Inches(0.85),
                fill=WHITE, line=col, corner=0.06, line_w=1.8)
        rect(s, Inches(0.6), y, Inches(0.12), Inches(0.85), fill=col)
        text(s, Inches(0.95), y + Inches(0.18), Inches(3.5), Inches(0.5),
             name, size=16, bold=True, color=col)
        chip(s, Inches(4.4), y + Inches(0.22), Inches(1.5), Inches(0.40),
             fields, fill=col, size=10)
        text(s, Inches(6.2), y + Inches(0.18), Inches(6.4), Inches(0.5),
             desc, size=12, color=TXT)

    # Footer note
    text(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.35),
         "M0 bepul tier · Singapore region · avtomatik backup · at-rest shifrlash",
         size=11, color=SUB, italic=True, align=PP_ALIGN.CENTER)
    footer(s, 10)


def s11_gamification(prs):
    s = blank(prs)
    header(s, "Gamifikatsiya formulalari",
           "Psixologik nazariyalarga asoslangan dizayn")

    # XP formula card
    rounded(s, Inches(0.6), Inches(1.85), Inches(5.9), Inches(2.2),
            fill=WHITE, line=PRIMARY, corner=0.06, line_w=2.0)
    text(s, Inches(0.85), Inches(2.0), Inches(5.4), Inches(0.45),
         "XP HISOBLASH", size=11, bold=True, color=PRIMARY)
    text(s, Inches(0.85), Inches(2.55), Inches(5.4), Inches(0.5),
         "XP = base × D_mult × SB",
         size=22, bold=True, color=INK)
    bullets(s, Inches(0.85), Inches(3.2), Inches(5.4), Inches(0.95), [
        "D_mult: easy 1.0 · medium 1.5 · hard 2.5 · expert 4.0",
        "SB(s) = min(1.5, 1 + 0.05·s) — streak bonusi",
    ], size=11)

    # Level formula card
    rounded(s, Inches(6.8), Inches(1.85), Inches(5.9), Inches(2.2),
            fill=WHITE, line=ACCENT, corner=0.06, line_w=2.0)
    text(s, Inches(7.05), Inches(2.0), Inches(5.4), Inches(0.45),
         "DARAJA TIZIMI", size=11, bold=True, color=ACCENT)
    text(s, Inches(7.05), Inches(2.55), Inches(5.4), Inches(0.5),
         "Level: 1 → 20 (eksponensial)",
         size=22, bold=True, color=INK)
    bullets(s, Inches(7.05), Inches(3.2), Inches(5.4), Inches(0.95), [
        "Past darajalarda kam XP — tez progress hissi",
        "Yuqori darajalarda ko'p XP — maqsad balandligi",
    ], size=11)

    # Streak protection
    rounded(s, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.4),
            fill=TINT, corner=0.06)
    text(s, Inches(0.95), Inches(4.45), Inches(11.5), Inches(0.45),
         "STREAK SAQLASH MEXANIZMI",
         size=12, bold=True, color=PRIMARY)
    bullets(s, Inches(0.95), Inches(4.95), Inches(11.5), Inches(1.7), [
        "Har kuni minimum 1 vazifa — streak davom etadi",
        "Streak Freeze: bir kun ko'tarib yuborilsa ham streak buzilmaydi",
        "10 kundan keyin SB maksimum 1.5× ga to'xtaydi — fairness saqlanadi",
        "Streak nolga tushganda re-engagement strategy aktivlashtiriladi",
    ], size=12)
    footer(s, 11)


def s12_api(prs):
    s = blank(prs)
    header(s, "RESTful API arxitekturasi",
           "33 endpoint · 6 router · P95 ≤ 300 ms")

    routers = [
        ("/auth",        "5 endpoint", "register · login · profile · password",        PRIMARY),
        ("/tasks",       "4 endpoint", "daily · recommended · complete · from-chat",  SECONDARY),
        ("/ai",          "5 endpoint", "chat · add-tasks · plan · insight · badges",  DANGER),
        ("/leaderboard", "3 endpoint", "global · weekly · user-rank",                 ACCENT),
        ("/progress",    "3 endpoint", "weekly · monthly · category breakdown",       SUCCESS),
        ("/users/me",    "13 endpoint","profile · data · export · sync · delete",    SKY),
    ]
    for i, (path, count, desc, col) in enumerate(routers):
        c, r = i % 3, i // 3
        x = Inches(0.6 + c * 4.15)
        y = Inches(1.85 + r * 2.2)
        rounded(s, x, y, Inches(3.95), Inches(2.0),
                fill=WHITE, line=col, corner=0.06, line_w=2.0)
        text(s, x + Inches(0.25), y + Inches(0.2), Inches(3.6), Inches(0.55),
             path, size=20, bold=True, color=col)
        chip(s, x + Inches(0.25), y + Inches(0.85), Inches(1.5), Inches(0.4),
             count, fill=col, size=10)
        text(s, x + Inches(0.25), y + Inches(1.35), Inches(3.6), Inches(0.6),
             desc, size=11, color=TXT)

    # bottom metrics
    text(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.4),
         "Auth: JWT Bearer (HMAC-SHA256) · Validatsiya: Pydantic v2 · Rate limit: 60/min",
         size=11, color=SUB, italic=True, align=PP_ALIGN.CENTER)
    footer(s, 12)


def s13_deploy(prs):
    s = blank(prs)
    header(s, "Joylashtirish (CI/CD)",
           "Avtomatik: git push → live deploy")

    pipeline = [
        ("Developer",  "Flutter + FastAPI\nlokal",       PRIMARY),
        ("GitHub",     "Push + Actions\nTest + Review",  INK),
        ("Render.com", "Webhook deploy\nHTTPS auto",     SUCCESS),
        ("MongoDB",    "Atlas Singapore\nM0 · backup",   SECONDARY),
    ]
    y = Inches(2.7)
    w = Inches(2.6)
    pad = Inches(0.35)
    x0 = Inches(0.7)
    for i, (name, sub, col) in enumerate(pipeline):
        x = x0 + (w + pad) * i
        rounded(s, x, y, w, Inches(2.0), fill=WHITE, line=col, corner=0.06, line_w=2.4)
        rect(s, x, y, w, Inches(0.55), fill=col)
        text(s, x, y + Inches(0.12), w, Inches(0.4),
             name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.2), y + Inches(0.85), w - Inches(0.4), Inches(1.0),
             sub, size=12, color=TXT, align=PP_ALIGN.CENTER)
        if i < len(pipeline) - 1:
            arrow(s, x + w, y + Inches(1.0), x + w + pad, y + Inches(1.0),
                  color=SUB, weight=2.5)
            label = ["git push", "webhook", "TCP/TLS"][i]
            text(s, x + w, y + Inches(1.1), pad, Inches(0.3),
                 label, size=9, color=SUB, italic=True, align=PP_ALIGN.CENTER)

    # Users + AI providers below
    rounded(s, Inches(0.7), Inches(5.5), Inches(5.9), Inches(1.2),
            fill=WHITE, line=ACCENT, corner=0.06, line_w=1.6)
    text(s, Inches(0.95), Inches(5.65), Inches(5.4), Inches(0.4),
         "FOYDALANUVCHILAR", size=11, bold=True, color=ACCENT)
    text(s, Inches(0.95), Inches(6.0), Inches(5.4), Inches(0.6),
         "iOS + Android · HTTPS REST API + JWT", size=12, color=TXT)

    rounded(s, Inches(6.8), Inches(5.5), Inches(5.9), Inches(1.2),
            fill=WHITE, line=DANGER, corner=0.06, line_w=1.6)
    text(s, Inches(7.05), Inches(5.65), Inches(5.4), Inches(0.4),
         "AI PROVIDERLAR", size=11, bold=True, color=DANGER)
    text(s, Inches(7.05), Inches(6.0), Inches(5.4), Inches(0.6),
         "OpenAI + Gemini + Groq · multi-fallback", size=12, color=TXT)
    footer(s, 13)


def s14_security(prs):
    s = blank(prs)
    header(s, "Xavfsizlik me'morchiligi", "Defense-in-depth — 6 qatlam")

    layers = [
        ("Transport",      "TLS 1.3 · Let's Encrypt",          DANGER),
        ("Auth",           "JWT Bearer · HMAC-SHA256",         ACCENT),
        ("Parol",          "bcrypt · work factor 12",          SUCCESS),
        ("Validatsiya",    "Pydantic v2 · type-safe",          SECONDARY),
        ("Database",       "Atlas IP whitelist · at-rest enc", PRIMARY),
        ("Rate limit",     "60/min · AI 30/h",                 SKY),
    ]
    for i, (title, desc, col) in enumerate(layers):
        c, r = i % 3, i // 3
        x = Inches(0.6 + c * 4.15)
        y = Inches(1.85 + r * 2.2)
        rounded(s, x, y, Inches(3.95), Inches(2.0),
                fill=WHITE, line=col, corner=0.06, line_w=2.0)
        circle(s, x + Inches(0.25), y + Inches(0.25), Inches(0.65), fill=col)
        text(s, x + Inches(0.25), y + Inches(0.32), Inches(0.65), Inches(0.5),
             str(i+1), size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, x + Inches(1.05), y + Inches(0.3), Inches(2.8), Inches(0.5),
             title, size=15, bold=True, color=INK)
        text(s, x + Inches(0.25), y + Inches(1.1), Inches(3.5), Inches(0.85),
             desc, size=12, color=TXT)

    text(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.35),
         "OWASP Top 10 zaifliklari ro'yxati to'liq qoplangan",
         size=11, color=SUB, italic=True, align=PP_ALIGN.CENTER)
    footer(s, 14)


def s15_results_kpi(prs):
    s = blank(prs)
    header(s, "Sinov natijalari · KPI",
           "n = 15 ishtirokchi · 7 kunlik foydalanish")

    kpis = [
        ("SUS",       "79.4 / 100", "Standart ≥ 68",  "Yaxshi",  SUCCESS),
        ("NPS",       "+42",        "Standart ≥ 30",  "Yaxshi",  SUCCESS),
        ("Bajarish",  "67 %",       "Standart ≥ 50%", "Yaxshi",  SUCCESS),
        ("AI chat",   "8.1 / 10",   "Standart ≥ 7",   "Yaxshi",  SUCCESS),
        ("Kunlik kirish", "3.8 / kun", "Standart ≥ 2", "A'lo",  PRIMARY),
        ("Cold start","35-55 sek",  "Ideal < 5 sek",  "Cheklov", DANGER),
    ]
    for i, (name, val, ref, badge, col) in enumerate(kpis):
        c, r = i % 3, i // 3
        x = Inches(0.6 + c * 4.15)
        y = Inches(1.85 + r * 2.2)
        rounded(s, x, y, Inches(3.95), Inches(2.0),
                fill=WHITE, line=col, corner=0.06, line_w=2.0)
        text(s, x + Inches(0.25), y + Inches(0.2), Inches(3.6), Inches(0.4),
             name, size=12, color=SUB)
        text(s, x + Inches(0.25), y + Inches(0.55), Inches(3.6), Inches(0.7),
             val, size=28, bold=True, color=col)
        text(s, x + Inches(0.25), y + Inches(1.25), Inches(3.6), Inches(0.35),
             ref, size=11, color=SUB, italic=True)
        chip(s, x + Inches(2.7), y + Inches(0.25), Inches(1.1), Inches(0.4),
             badge, fill=col, size=10)
    footer(s, 15)


def s16_results_growth(prs):
    s = blank(prs)
    header(s, "Foydalanish dinamikasi", "7 kunlik kirish chastotasi")

    days_data = [2.1, 3.4, 3.9, 4.2, 3.8, 4.1, 3.8]
    labels = ["1-kun", "2-kun", "3-kun", "4-kun", "5-kun", "6-kun", "7-kun"]

    # Mini bar chart with shapes
    chart_x = Inches(1.0)
    chart_y = Inches(2.0)
    chart_w = Inches(11.3)
    chart_h = Inches(4.0)

    max_val = 5.0
    bar_w = Inches(1.0)
    gap = (chart_w - bar_w * 7) / 8  # 8 gaps including edges
    threshold_y = chart_y + chart_h * (1 - 2.0 / max_val)

    # baseline
    arrow(s, chart_x, chart_y + chart_h,
          chart_x + chart_w, chart_y + chart_h, color=SUB, weight=1.0)
    # standard line
    line = s.shapes.add_connector(2, chart_x, threshold_y,
                                  chart_x + chart_w, threshold_y)
    line.line.color.rgb = DANGER
    line.line.width = Pt(1.5)
    line.line.dash_style = 7  # DASH
    text(s, chart_x + chart_w - Inches(2.5), threshold_y - Inches(0.35),
         Inches(2.4), Inches(0.3),
         "Soha standarti (≥ 2/kun)", size=10,
         color=DANGER, italic=True, align=PP_ALIGN.RIGHT)

    for i, (val, lbl) in enumerate(zip(days_data, labels)):
        x = chart_x + gap + (bar_w + gap) * i
        h = chart_h * (val / max_val)
        y = chart_y + chart_h - h
        color = ACCENT if val < 2.5 else PRIMARY
        rect(s, x, y, bar_w, h, fill=color)
        # value label
        text(s, x, y - Inches(0.45), bar_w, Inches(0.4),
             f"{val}", size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        # x label
        text(s, x, chart_y + chart_h + Inches(0.1), bar_w, Inches(0.3),
             lbl, size=10, color=SUB, align=PP_ALIGN.CENTER)

    # Bottom takeaway
    text(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
         "2-kundan boshlab standartdan ustun · 4-kunga 'odat shakllanishi' yetadi",
         size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    footer(s, 16)


def s17_advantages(prs):
    s = blank(prs)
    header(s, "MotivAI ning afzalliklari", "Nima uchun innovatsion?")

    items = [
        ("O'zbek tilida lokallashtirilgan",
         "700+ tarjima kaliti · 3 til real-time switch · mahalliy kontekst"),
        ("Multi-provider AI fallback",
         "OpenAI → Gemini → Groq · 99.6% uptime · quota cheklovi yo'q"),
        ("Real vaqtdagi shaxsiylashtirish",
         "MVF har so'rovda qayta hisoblab · profil dinamik yangilanadi"),
        ("Psixologik nazariyalar bazasi",
         "SDT + Flow Theory + Gamification matematik formalizatsiyasi"),
        ("Cross-platform yagona kod",
         "Flutter · iOS+Android · 64.8 MB APK · Skia native render"),
        ("Ochiq manba",
         "GitHub · MIT litsenziya · modulli kengaytiriladigan arxitektura"),
    ]
    for i, (title, desc) in enumerate(items):
        c, r = i % 2, i // 2
        x = Inches(0.6 + c * 6.15)
        y = Inches(1.85 + r * 1.65)
        rounded(s, x, y, Inches(5.95), Inches(1.5),
                fill=WHITE, line=PRIMARY, corner=0.06, line_w=1.6)
        # number badge
        circle(s, x + Inches(0.25), y + Inches(0.3), Inches(0.5), fill=PRIMARY)
        text(s, x + Inches(0.25), y + Inches(0.36), Inches(0.5), Inches(0.4),
             str(i+1), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.95), y + Inches(0.2), Inches(4.9), Inches(0.5),
             title, size=14, bold=True, color=INK)
        text(s, x + Inches(0.95), y + Inches(0.75), Inches(4.9), Inches(0.7),
             desc, size=11, color=TXT)
    footer(s, 17)


def s18_innovations(prs):
    s = blank(prs)
    header(s, "Ilmiy yangiliklar va hissalar",
           "Loyihaning original ko'rsatkichlari")

    rounded(s, Inches(0.6), Inches(1.85), Inches(12.1), Inches(1.4),
            fill=PRIMARY, corner=0.06)
    text(s, Inches(0.85), Inches(2.0), Inches(11.5), Inches(0.45),
         "ASOSIY ILMIY YANGILIK",
         size=11, bold=True, color=RGBColor(0xE0, 0xE7, 0xFF))
    text(s, Inches(0.85), Inches(2.45), Inches(11.5), Inches(0.85),
         "O'zbek tilida ishlovchi, multi-provider AI fallback va psixologik "
         "gamifikatsiyani birlashtirgan to'liq funksional mobil platforma — "
         "O'zbekistonda birinchi marta yaratildi.",
         size=14, color=WHITE)

    items = [
        ("MVF gibrid model",   "NDCG@5 = 0.78", "Sanoat darajasidagi tavsiya sifati"),
        ("Multi-AI chain",     "3 provider",     "Yangi turdagi fallback arxitektura"),
        ("Empirik baholash",   "SUS = 79.4",     "Soha standartidan 16.8% ustun"),
        ("Lokallashtirish",    "700+ kalit",     "3 tilda to'liq native qo'llab-quvvatlash"),
    ]
    for i, (title, val, desc) in enumerate(items):
        c, r = i % 2, i // 2
        x = Inches(0.6 + c * 6.15)
        y = Inches(3.5 + r * 1.65)
        rounded(s, x, y, Inches(5.95), Inches(1.5),
                fill=WHITE, line=ACCENT, corner=0.06, line_w=1.6)
        text(s, x + Inches(0.3), y + Inches(0.18), Inches(2.4), Inches(0.5),
             title, size=13, bold=True, color=INK)
        text(s, x + Inches(0.3), y + Inches(0.75), Inches(2.4), Inches(0.5),
             val, size=22, bold=True, color=ACCENT)
        text(s, x + Inches(2.85), y + Inches(0.3), Inches(2.95), Inches(1.0),
             desc, size=12, color=TXT)
    footer(s, 18)


def s19_conclusions(prs):
    s = blank(prs)
    header(s, "Xulosa va kelajakdagi yo'nalishlar")

    # Two columns
    rounded(s, Inches(0.6), Inches(1.85), Inches(5.95), Inches(5.0),
            fill=WHITE, line=SUCCESS, corner=0.06, line_w=2.0)
    text(s, Inches(0.85), Inches(2.0), Inches(5.5), Inches(0.5),
         "✓  ERISHILGAN", size=14, bold=True, color=SUCCESS)
    bullets(s, Inches(0.85), Inches(2.55), Inches(5.5), Inches(4.2), [
        "33 endpoint · 94 ms o'rtacha javob",
        "MVF NDCG@5 = 0.78",
        "99.6% uptime (multi-provider AI)",
        "SUS = 79.4 · NPS = +42",
        "iOS+Android yagona kod bazasi",
        "Render bepul tier'da ishlaydi",
        "700+ tarjima kaliti · 3 til",
    ], size=12)

    rounded(s, Inches(6.75), Inches(1.85), Inches(5.95), Inches(5.0),
            fill=WHITE, line=ACCENT, corner=0.06, line_w=2.0)
    text(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.5),
         "→  KELAJAK", size=14, bold=True, color=ACCENT)
    bullets(s, Inches(7.0), Inches(2.55), Inches(5.5), Inches(4.2), [
        "BERT4Rec transformer (+8-12% NDCG)",
        "Hive/Isar — offline rejim",
        "FCM push bildirishnomalar",
        "Redis kesh — AI 60-70% tezroq",
        "HEMIS integratsiya (universitet)",
        "A/B testlash platforma",
        "300+ talaba · 6 oylik kengaytirilgan",
    ], size=12)
    footer(s, 19)


def s20_thanks(prs):
    s = blank(prs)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, fill=PRIMARY)

    text(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.0),
         "RAHMAT!", size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.5),
         "Savol va takliflar uchun ochiqman",
         size=20, color=RGBColor(0xE0, 0xE7, 0xFF), italic=True,
         align=PP_ALIGN.CENTER)

    rounded(s, Inches(4.0), Inches(4.75), Inches(5.3), Inches(1.6),
            fill=WHITE, corner=0.10)
    text(s, Inches(4.0), Inches(4.90), Inches(5.3), Inches(0.4),
         "Abduvaliyev Samandar", size=14, bold=True, color=INK,
         align=PP_ALIGN.CENTER)
    text(s, Inches(4.0), Inches(5.35), Inches(5.3), Inches(0.4),
         "github.com/AbduvaliyevSamandar/motivai", size=12,
         color=PRIMARY, align=PP_ALIGN.CENTER)
    text(s, Inches(4.0), Inches(5.80), Inches(5.3), Inches(0.4),
         "elmurodovmaxmud77@gmail.com", size=12, color=SUB,
         align=PP_ALIGN.CENTER)


# ── main ───────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    setup(prs)
    builders = [
        s01_title, s02_problem, s03_goal_tasks, s04_architecture_mini,
        s05_techstack, s06_archetypes, s07_mvf, s08_user_flow,
        s09_ai_chain, s10_database, s11_gamification, s12_api,
        s13_deploy, s14_security, s15_results_kpi, s16_results_growth,
        s17_advantages, s18_innovations, s19_conclusions, s20_thanks,
    ]
    for fn in builders:
        fn(prs)
    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
